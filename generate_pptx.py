#!/usr/bin/env python3
"""Generate a 14-slide BlockTrace investor pitch deck as PPTX."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Dimensions
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Colours
BG_DARK = RGBColor(0x0A, 0x0F, 0x1C)
CARD_BG = RGBColor(0x1E, 0x29, 0x3B)
DARKER_BG = RGBColor(0x0F, 0x17, 0x2A)
ACCENT = RGBColor(0x22, 0xD3, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LIGHT = RGBColor(0xB4, 0xBF, 0xCC)
GRAY_MED = RGBColor(0x94, 0xA3, 0xB8)
GRAY_DARK = RGBColor(0x72, 0x81, 0x97)
BLACK = RGBColor(0x0A, 0x0F, 0x1C)


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, font_name="Inter",
                 alignment=PP_ALIGN.LEFT, line_spacing=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    return txBox


def add_rounded_rect(slide, left, top, width, height,
                     fill_color=CARD_BG, corner=Inches(0.15)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if hasattr(shape, "adjustments") and len(shape.adjustments) > 0:
        shape.adjustments[0] = corner / width if width > 0 else 0.05
    return shape


def add_accent_line(slide, left, top, width, height=Pt(3)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape


def add_slide_number(slide, number):
    add_text_box(
        slide, Inches(12.4), Inches(6.95), Inches(0.8), Inches(0.4),
        f"{number:02d}", font_size=10, color=GRAY_DARK,
        font_name="JetBrains Mono", alignment=PP_ALIGN.RIGHT
    )


def add_section_label(slide, text, left=Inches(0.83), top=Inches(0.55)):
    add_text_box(
        slide, left, top, Inches(4), Inches(0.3),
        text, font_size=10, color=ACCENT, bold=True,
        font_name="JetBrains Mono"
    )


def add_title(slide, text, left=Inches(0.83), top=Inches(0.85),
              width=Inches(8), font_size=36):
    add_text_box(
        slide, left, top, width, Inches(0.8),
        text, font_size=font_size, color=WHITE, bold=True
    )


def add_description(slide, text, left=Inches(0.83), top=Inches(1.55),
                    width=Inches(5), font_size=14):
    add_text_box(
        slide, left, top, width, Inches(1.2),
        text, font_size=font_size, color=GRAY_MED, line_spacing=22
    )


def add_card(slide, left, top, width, height, title, desc,
             icon_text="*", title_size=14, desc_size=11):
    add_rounded_rect(slide, left, top, width, height)
    add_text_box(
        slide, left + Inches(0.25), top + Inches(0.25),
        Inches(0.4), Inches(0.35),
        icon_text, font_size=18, color=ACCENT, bold=True
    )
    add_text_box(
        slide, left + Inches(0.25), top + Inches(0.6),
        width - Inches(0.5), Inches(0.35),
        title, font_size=title_size, color=WHITE, bold=True
    )
    add_text_box(
        slide, left + Inches(0.25), top + Inches(0.95),
        width - Inches(0.5), height - Inches(1.2),
        desc, font_size=desc_size, color=GRAY_LIGHT, line_spacing=18
    )


def add_card_brickwork(slide, left, top, width, height, title, desc,
                       icon_text="*", badge_text=None, badge_accent=False,
                       title_size=16, desc_size=11):
    add_rounded_rect(slide, left, top, width, height)
    y_cursor = top + Inches(0.25)
    if badge_text:
        bw, bh = Inches(1.2), Inches(0.22)
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.25), y_cursor, bw, bh
        )
        if badge_accent:
            badge.fill.solid()
            badge.fill.fore_color.rgb = ACCENT
            badge_color = BLACK
        else:
            badge.fill.solid()
            badge.fill.fore_color.rgb = DARKER_BG
            badge_color = ACCENT
        badge.line.fill.background()
        add_text_box(
            slide, left + Inches(0.25), y_cursor, bw, bh, badge_text,
            font_size=8, color=badge_color, bold=True,
            font_name="JetBrains Mono", alignment=PP_ALIGN.CENTER
        )
        y_cursor += Inches(0.35)
    add_text_box(
        slide, left + Inches(0.25), y_cursor,
        Inches(0.4), Inches(0.35),
        icon_text, font_size=18, color=ACCENT, bold=True
    )
    y_cursor += Inches(0.35)
    add_text_box(
        slide, left + Inches(0.25), y_cursor,
        width - Inches(0.5), Inches(0.35),
        title, font_size=title_size, color=WHITE, bold=True
    )
    y_cursor += Inches(0.35)
    add_text_box(
        slide, left + Inches(0.25), y_cursor,
        width - Inches(0.5), height - (y_cursor - top) - Inches(0.25),
        desc, font_size=desc_size, color=GRAY_LIGHT, line_spacing=18
    )


def add_tree_node(slide, left, top, width, height, text,
                  fill=CARD_BG, text_color=WHITE, font_size=12,
                  border_color=None, bold=False):
    shape = add_rounded_rect(slide, left, top, width, height, fill_color=fill)
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    add_text_box(
        slide, left, top, width, height,
        text, font_size=font_size, color=text_color, bold=bold,
        alignment=PP_ALIGN.CENTER
    )


def add_connector_line(slide, left, top, width, height, color=ACCENT):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


# ── Build Presentation ─────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 01 - Title
# ═══════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
set_slide_bg(s1)
add_accent_line(s1, Inches(0), Inches(0), SLIDE_W, Pt(3))
add_text_box(s1, Inches(2), Inches(2.2), Inches(9.333), Inches(0.9),
             "BlockTrace", font_size=48, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(s1, Inches(3.5), Inches(3.1), Inches(6.333), Inches(0.6),
             "Composable Token Infrastructure for Real-World Assets",
             font_size=20, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_accent_line(s1, Inches(6), Inches(3.8), Inches(1.333), Pt(2))
add_text_box(s1, Inches(3), Inches(4.1), Inches(7.333), Inches(1),
             "Turn physical assets into verifiable digital structures with "
             "complete lifecycle traceability.",
             font_size=14, color=GRAY_LIGHT, alignment=PP_ALIGN.CENTER,
             line_spacing=22)
add_slide_number(s1, 1)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 02 - The Problem
# ═══════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
set_slide_bg(s2)
add_section_label(s2, "THE PROBLEM")
add_title(s2, "Traceability Is Broken")
add_description(s2,
    "Enterprise asset data lives in fragmented silos. Compliance is manual. "
    "Cross-organisation trust is non-existent. Bills of materials are "
    "unverifiable.")

problem_cards = [
    ("*", "Siloed Databases",
     "Siloed enterprise databases with no interoperability. Traceability data "
     "lives across dozens of disconnected ERPs, spreadsheets, and legacy "
     "databases."),
    ("*", "No Cross-Org Trust",
     "No cross-organisation trust layer. Partners, auditors, and regulators "
     "cannot independently verify claims about asset provenance."),
    ("*", "Unverifiable BOMs",
     "Non-verifiable bills of materials. Component history is easily lost or "
     "falsified across supply chain handoffs."),
    ("*", "Manual Compliance",
     "Manual compliance and audit processes. Regulatory compliance is handled "
     "through costly, unscalable manual processes."),
    ("*", "Fragmented Lifecycle",
     "Fragmented asset lifecycle records. No single source of truth for an "
     "asset's complete history across organisations."),
]

card_w = Inches(2.35)
card_h = Inches(2.8)
x_start = Inches(0.83)
y_cards = Inches(3.0)
card_gap = Inches(0.2)

for i, (icon, title, desc) in enumerate(problem_cards):
    x = x_start + i * (card_w + card_gap)
    add_card(s2, x, y_cards, card_w, card_h, title, desc, icon_text=icon)
add_slide_number(s2, 2)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 03 - The Opportunity
# ═══════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
set_slide_bg(s3)
add_section_label(s3, "THE OPPORTUNITY")
add_title(s3, "Three Forces Converging")
add_description(s3,
    "Regulatory pressure, enterprise SaaS maturity, and RWA tokenisation "
    "are creating an infrastructure-level opportunity.", width=Inches(5.5))

pill_w = Inches(2.8)
pill_h = Inches(1.6)
pill_y = Inches(3.0)
pills = [
    ("Blockchain\nInfrastructure",),
    ("Verifiable\nAsset Infrastructure",),
    ("Enterprise SaaS\nMaturity",),
]
pill_positions = [Inches(0.83), Inches(5.27), Inches(9.7)]

for i in range(3):
    text = pills[i][0]
    px = pill_positions[i]
    is_center = (i == 1)
    h = pill_h + Inches(0.3) if is_center else pill_h
    y = pill_y - Inches(0.15) if is_center else pill_y
    shape = add_rounded_rect(s3, px, y, pill_w, h, fill_color=CARD_BG)
    if is_center:
        shape.line.color.rgb = ACCENT
        shape.line.width = Pt(2)
    add_text_box(s3, px + Inches(0.3), y + Inches(0.4),
                 pill_w - Inches(0.6), Inches(0.8),
                 text, font_size=14 if not is_center else 16,
                 color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

for ax in [Inches(3.63), Inches(8.07)]:
    add_text_box(s3, ax, pill_y + Inches(0.5), Inches(1.64), Inches(0.5),
                 "\u2192", font_size=24, color=ACCENT, bold=True,
                 alignment=PP_ALIGN.CENTER)

stats = [
    ("$16T+", "Projected RWA tokenisation market by 2030"),
    ("80%", "Of enterprises cite supply chain visibility as critical priority"),
    ("47+", "New supply chain due diligence regulations enacted globally "
     "since 2020"),
]
stat_w = Inches(3.5)
stat_y = Inches(5.2)
for i, (val, label) in enumerate(stats):
    sx = Inches(0.83) + i * (stat_w + Inches(0.65))
    add_text_box(s3, sx, stat_y, stat_w, Inches(0.6),
                 val, font_size=30, color=ACCENT, bold=True,
                 font_name="JetBrains Mono")
    add_text_box(s3, sx, stat_y + Inches(0.6), stat_w, Inches(0.6),
                 label, font_size=11, color=GRAY_LIGHT, line_spacing=17)
add_slide_number(s3, 3)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 04 - The Solution (Tree Diagrams)
# ═══════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
set_slide_bg(s4)
add_section_label(s4, "THE SOLUTION")
add_title(s4, "From Flat Records to Composable Asset Graphs", width=Inches(9))
add_description(s4,
    "A real-world asset becomes a root token. Each certificate, component, "
    "or document is a sub-token. Each sub-token can itself contain "
    "sub-tokens \u2014 forming a verifiable, composable structure.",
    width=Inches(5))

# Real Estate Example
re_x = Inches(0.83)
re_y = Inches(3.0)
add_text_box(s4, re_x, re_y - Inches(0.3), Inches(3), Inches(0.25),
             "REAL ESTATE EXAMPLE", font_size=9, color=GRAY_DARK,
             bold=True, font_name="JetBrains Mono")

root_w, root_h = Inches(2.6), Inches(0.45)
root_x = re_x + Inches(1.3)
add_tree_node(s4, root_x, re_y, root_w, root_h,
              "House Token", fill=ACCENT, text_color=BLACK,
              font_size=13, bold=True)

trunk_x = root_x + root_w // 2
add_connector_line(s4, trunk_x, re_y + root_h, Pt(2), Inches(0.3))

children_re = ["Survey Token", "Electrical Certificate",
               "Renovation Record", "Insurance Policy"]
child_w, child_h, child_gap = Inches(1.4), Inches(0.38), Inches(0.12)
total_children_w = len(children_re) * child_w + (len(children_re) - 1) * child_gap
children_start_x = root_x + root_w // 2 - total_children_w // 2
hbar_y = re_y + root_h + Inches(0.3)
add_connector_line(s4, children_start_x + child_w // 2, hbar_y,
                   total_children_w - child_w, Pt(2))

child_y = hbar_y + Inches(0.2)
for i, name in enumerate(children_re):
    cx = children_start_x + i * (child_w + child_gap)
    add_connector_line(s4, cx + child_w // 2, hbar_y, Pt(2), Inches(0.2))
    add_tree_node(s4, cx, child_y, child_w, child_h, name,
                  fill=CARD_BG, text_color=WHITE, font_size=10,
                  border_color=ACCENT)

sub_trunk_x = children_start_x + child_w // 2
add_connector_line(s4, sub_trunk_x, child_y + child_h, Pt(2), Inches(0.35))

sub_names = ["Site Plan", "Boundary Report"]
sub_w, sub_h, sub_gap = Inches(1.1), Inches(0.32), Inches(0.1)
total_sub_w = len(sub_names) * sub_w + (len(sub_names) - 1) * sub_gap
sub_start_x = sub_trunk_x - total_sub_w // 2
sub_hbar_y = child_y + child_h + Inches(0.3)
add_connector_line(s4, sub_start_x + sub_w // 2, sub_hbar_y,
                   total_sub_w - sub_w, Pt(2))

sub_y = child_y + child_h + Inches(0.5)
for i, sn in enumerate(sub_names):
    sx = sub_start_x + i * (sub_w + sub_gap)
    add_connector_line(s4, sx + sub_w // 2, sub_hbar_y, Pt(2), Inches(0.18))
    add_tree_node(s4, sx, sub_y, sub_w, sub_h, sn,
                  fill=DARKER_BG, text_color=GRAY_LIGHT, font_size=9,
                  border_color=ACCENT)

# Manufacturing Example
mfg_x = Inches(7)
mfg_y = Inches(3.0)
add_text_box(s4, mfg_x, mfg_y - Inches(0.3), Inches(3), Inches(0.25),
             "MANUFACTURING EXAMPLE", font_size=9, color=GRAY_DARK,
             bold=True, font_name="JetBrains Mono")

m_root_x = mfg_x + Inches(1.3)
add_tree_node(s4, m_root_x, mfg_y, root_w, root_h,
              "Finished Product", fill=CARD_BG, text_color=WHITE,
              font_size=13, bold=True, border_color=ACCENT)

add_connector_line(s4, m_root_x + root_w // 2, mfg_y + root_h,
                   Pt(2), Inches(0.3))

children_mfg = ["Sub-Assembly A", "Quality Certificate", "Shipping Manifest"]
m_child_w, m_child_gap = Inches(1.55), Inches(0.15)
m_total_w = len(children_mfg) * m_child_w + (len(children_mfg) - 1) * m_child_gap
m_children_start = m_root_x + root_w // 2 - m_total_w // 2
m_hbar_y = mfg_y + root_h + Inches(0.3)
add_connector_line(s4, m_children_start + m_child_w // 2, m_hbar_y,
                   m_total_w - m_child_w, Pt(2))

m_child_y = m_hbar_y + Inches(0.2)
for i, name in enumerate(children_mfg):
    cx = m_children_start + i * (m_child_w + m_child_gap)
    add_connector_line(s4, cx + m_child_w // 2, m_hbar_y, Pt(2), Inches(0.2))
    add_tree_node(s4, cx, m_child_y, m_child_w, child_h, name,
                  fill=CARD_BG, text_color=WHITE, font_size=10,
                  border_color=ACCENT)

# Sub-Assembly A sub-tokens
sa_trunk_x = m_children_start + m_child_w // 2
add_connector_line(s4, sa_trunk_x, m_child_y + child_h, Pt(2), Inches(0.35))
sa_subs = ["Component A1", "Component A2"]
sa_sub_hbar_y = m_child_y + child_h + Inches(0.3)
sa_sub_start = sa_trunk_x - total_sub_w // 2
add_connector_line(s4, sa_sub_start + sub_w // 2, sa_sub_hbar_y,
                   total_sub_w - sub_w, Pt(2))
sa_sub_y = m_child_y + child_h + Inches(0.5)
for i, sn in enumerate(sa_subs):
    sx = sa_sub_start + i * (sub_w + sub_gap)
    add_connector_line(s4, sx + sub_w // 2, sa_sub_hbar_y, Pt(2), Inches(0.18))
    add_tree_node(s4, sx, sa_sub_y, sub_w, sub_h, sn,
                  fill=DARKER_BG, text_color=GRAY_LIGHT, font_size=9,
                  border_color=ACCENT)

# Quality Certificate sub-tokens
qc_trunk_x = m_children_start + m_child_w + m_child_gap + m_child_w // 2
add_connector_line(s4, qc_trunk_x, m_child_y + child_h, Pt(2), Inches(0.35))
qc_subs = ["Lab Test Report", "Compliance Cert"]
qc_sub_start = qc_trunk_x - total_sub_w // 2
add_connector_line(s4, qc_sub_start + sub_w // 2, sa_sub_hbar_y,
                   total_sub_w - sub_w, Pt(2))
for i, sn in enumerate(qc_subs):
    sx = qc_sub_start + i * (sub_w + sub_gap)
    add_connector_line(s4, sx + sub_w // 2, sa_sub_hbar_y, Pt(2), Inches(0.18))
    add_tree_node(s4, sx, sa_sub_y, sub_w, sub_h, sn,
                  fill=DARKER_BG, text_color=GRAY_LIGHT, font_size=9,
                  border_color=ACCENT)
add_slide_number(s4, 4)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 05 - How It Works
# ═══════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank_layout)
set_slide_bg(s5)
add_section_label(s5, "HOW IT WORKS")
add_title(s5, "A Layered Architecture Built for Enterprise", width=Inches(10))
add_description(s5,
    "Five purpose-built layers work together to tokenise, anchor, and query "
    "your asset data \u2014 without replacing existing systems.",
    width=Inches(5.5))

layers = [
    ("1", "Integration Layer",
     "REST APIs, ERP connectors, webhooks, batch import. Enterprise systems "
     "push asset data in.",
     ["REST API", "ERP", "Webhooks"], True),
    ("2", "Tokenisation Engine",
     "Hierarchical token pack creation, versioning, composition. Schema "
     "validation and lifecycle event tracking.",
     ["Core"], True),
    ("3", "Off-Chain Indexed Data Layer",
     "Structured storage, fast retrieval, selective disclosure. Full asset "
     "graph traversal.",
     ["Indexed"], False),
    ("4", "On-Chain Hash Anchoring",
     "Immutable proof on any blockchain, chain-agnostic. Tamper-evident, "
     "cryptographic audit trail.",
     ["Chain-Agnostic"], False),
    ("5", "Query & Analytics Dashboard",
     "Real-time asset insights, dependency maps, lifecycle views. Compliance "
     "reporting and event monitoring.",
     ["Dashboard"], False),
]

layer_w = Inches(10.5)
layer_h = Inches(0.75)
layer_x = Inches(1.4)
layer_start_y = Inches(2.8)
layer_gap = Inches(0.12)

for i, (num, ltitle, ldesc, tags, highlight) in enumerate(layers):
    ly = layer_start_y + i * (layer_h + layer_gap)
    shape = add_rounded_rect(s5, layer_x, ly, layer_w, layer_h,
                             fill_color=CARD_BG)
    if highlight and i == 1:
        shape.line.color.rgb = ACCENT
        shape.line.width = Pt(1)

    badge_sz = Inches(0.35)
    add_rounded_rect(s5, layer_x + Inches(0.3),
                     ly + (layer_h - badge_sz) / 2,
                     badge_sz, badge_sz, fill_color=ACCENT)
    add_text_box(s5, layer_x + Inches(0.3),
                 ly + (layer_h - badge_sz) / 2,
                 badge_sz, badge_sz,
                 num, font_size=14, color=BLACK, bold=True,
                 font_name="JetBrains Mono", alignment=PP_ALIGN.CENTER)
    add_text_box(s5, layer_x + Inches(0.8), ly + Inches(0.08),
                 Inches(3), Inches(0.3),
                 ltitle, font_size=14, color=WHITE, bold=True)
    add_text_box(s5, layer_x + Inches(0.8), ly + Inches(0.38),
                 Inches(6), Inches(0.35),
                 ldesc, font_size=11, color=GRAY_LIGHT)

    tag_x = layer_x + Inches(8.5)
    for j, tag in enumerate(tags):
        tw = Inches(1.0)
        tx = tag_x + j * Inches(1.1)
        add_rounded_rect(s5, tx, ly + Inches(0.22), tw, Inches(0.3),
                         fill_color=DARKER_BG)
        add_text_box(s5, tx, ly + Inches(0.22), tw, Inches(0.3),
                     tag, font_size=9, color=ACCENT,
                     font_name="JetBrains Mono",
                     alignment=PP_ALIGN.CENTER)
add_slide_number(s5, 5)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 06 - Why Token Packs Matter (Brickwork)
# ═══════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank_layout)
set_slide_bg(s6)
add_section_label(s6, "WHY IT MATTERS")
add_title(s6, "Why Hierarchical Token Packs Matter", width=Inches(8))
add_description(s6,
    "Hierarchical Token Packs enable capabilities that flat tokenisation "
    "systems simply cannot deliver.", width=Inches(5.5))

brick_w = Inches(3.67)
brick_h = Inches(2.5)
brick_gap = Inches(0.2)
brick_y_top = Inches(3.0)
brick_y_bot = brick_y_top + brick_h + brick_gap

bricks6_top = [
    ("*", "Recursive Provenance",
     "Trace any component back through its full history, across every "
     "level of assembly."),
    ("*", "Component-Level Recall",
     "Identify and isolate affected assets instantly when a component "
     "is recalled."),
    ("*", "Full Lifecycle Versioning",
     "Every change to an asset or sub-token is versioned, timestamped, "
     "and immutable."),
]
bricks6_bot = [
    ("*", "Verifiable BOM",
     "Cryptographically verify every bill of materials down to the "
     "component level."),
    ("*", "Audit-Ready Assets",
     "Generate compliance reports instantly with cryptographic proof of "
     "every claim."),
]

for i, (icon, t, d) in enumerate(bricks6_top):
    bx = Inches(0.83) + i * (brick_w + brick_gap)
    add_card(s6, bx, brick_y_top, brick_w, brick_h, t, d, icon_text=icon)

bot_total = 2 * brick_w + brick_gap
bot_start = Inches(0.83) + (3 * brick_w + 2 * brick_gap - bot_total) / 2
for i, (icon, t, d) in enumerate(bricks6_bot):
    bx = bot_start + i * (brick_w + brick_gap)
    add_card(s6, bx, brick_y_bot, brick_w, brick_h, t, d, icon_text=icon)
add_slide_number(s6, 6)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 07 - Use Cases
# ═══════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank_layout)
set_slide_bg(s7)
add_section_label(s7, "USE CASES")
add_title(s7, "Built for Asset-Heavy Industries")

use_cases = [
    ("Manufacturing", True,
     "Track sub-assemblies, components, and certifications across complex "
     "production lines.",
     [("Engine Assembly", ACCENT), ("\u251C\u2500 Cylinder Block", GRAY_LIGHT),
      ("\u251C\u2500 Crankshaft", GRAY_LIGHT),
      ("\u2514\u2500 QC Report", GRAY_MED)]),
    ("Real Estate", False,
     "Compose property tokens from surveys, certificates, and renovation "
     "records.",
     [("Property Token", ACCENT), ("\u251C\u2500 Title Deed", GRAY_LIGHT),
      ("\u251C\u2500 Survey Report", GRAY_LIGHT),
      ("\u2514\u2500 Energy Rating", GRAY_MED)]),
    ("Pharmaceuticals", False,
     "Full chain-of-custody from raw material to patient delivery with "
     "regulatory compliance.",
     [("Drug Batch", ACCENT), ("\u251C\u2500 Raw Material Cert", GRAY_LIGHT),
      ("\u251C\u2500 Lab Analysis", GRAY_LIGHT),
      ("\u2514\u2500 Chain of Custody", GRAY_MED)]),
    ("Energy & Carbon", False,
     "Verifiable carbon credits and energy asset provenance with embedded "
     "compliance.",
     [("Carbon Credit", ACCENT), ("\u251C\u2500 Project Audit", GRAY_LIGHT),
      ("\u251C\u2500 Measurement Data", GRAY_LIGHT),
      ("\u2514\u2500 Retirement Record", GRAY_MED)]),
]

uc_w = Inches(2.85)
uc_gap = Inches(0.2)
uc_x = Inches(0.83)
uc_y = Inches(2.0)

for i, (name, is_primary, desc, tree) in enumerate(use_cases):
    cx = uc_x + i * (uc_w + uc_gap)
    header_h = Inches(0.45)
    hdr_fill = ACCENT if is_primary else DARKER_BG
    hdr_color = BLACK if is_primary else WHITE
    add_rounded_rect(s7, cx, uc_y, uc_w, header_h, fill_color=hdr_fill)
    add_text_box(s7, cx + Inches(0.2), uc_y, uc_w - Inches(0.4), header_h,
                 name, font_size=13, color=hdr_color, bold=True)

    body_h = Inches(4.5)
    add_rounded_rect(s7, cx, uc_y + header_h, uc_w, body_h, fill_color=CARD_BG)
    add_text_box(s7, cx + Inches(0.2), uc_y + header_h + Inches(0.15),
                 uc_w - Inches(0.4), Inches(0.9),
                 desc, font_size=11, color=GRAY_LIGHT, line_spacing=17)

    tree_y = uc_y + header_h + Inches(1.1)
    tree_h = Inches(1.8)
    add_rounded_rect(s7, cx + Inches(0.15), tree_y,
                     uc_w - Inches(0.3), tree_h, fill_color=DARKER_BG)
    for j, (line_text, line_color) in enumerate(tree):
        add_text_box(s7, cx + Inches(0.3),
                     tree_y + Inches(0.15) + j * Inches(0.35),
                     uc_w - Inches(0.6), Inches(0.3),
                     line_text, font_size=10, color=line_color,
                     font_name="JetBrains Mono", bold=(j == 0))
add_slide_number(s7, 7)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 08 - Why Integrate (Brickwork)
# ═══════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(blank_layout)
set_slide_bg(s8)
add_section_label(s8, "INTEGRATION")
add_title(s8, "Already Have Traceability? Even Better.", width=Inches(8))
add_description(s8,
    "BlockTrace doesn\u2019t replace your existing systems \u2014 it makes them "
    "provable, composable, and trusted across your entire supply chain.",
    width=Inches(6.5))

int8_top = [
    ("*", "Complement, Not Replace",
     "Works alongside SAP, Oracle, or custom ERP systems. No "
     "rip-and-replace \u2014 just a trust layer on top of what you "
     "already have."),
    ("*", "Cryptographic Proof Layer",
     "Add tamper-evident, hash-verified proof on top of your existing "
     "records. Anchor to any blockchain without changing your workflow."),
    ("*", "API-First Integration",
     "REST APIs, webhooks, ERP connectors, and batch import. Connect "
     "your existing systems in days, not months."),
]
int8_bot = [
    ("*", "Cross-Organisation Trust",
     "Your internal system tracks your data. BlockTrace proves it to "
     "partners, regulators, and customers \u2014 without exposing "
     "sensitive details."),
    ("*", "From Flat to Composable",
     "Transform siloed, flat records into hierarchical token packs with "
     "recursive provenance \u2014 turning your data into a verifiable "
     "asset graph."),
]

for i, (icon, t, d) in enumerate(int8_top):
    bx = Inches(0.83) + i * (brick_w + brick_gap)
    add_card(s8, bx, brick_y_top, brick_w, brick_h, t, d, icon_text=icon)

for i, (icon, t, d) in enumerate(int8_bot):
    bx = bot_start + i * (brick_w + brick_gap)
    add_card(s8, bx, brick_y_bot, brick_w, brick_h, t, d, icon_text=icon)
add_slide_number(s8, 8)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 09 - Competitive Landscape
# ═══════════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(blank_layout)
set_slide_bg(s9)
add_section_label(s9, "COMPETITIVE LANDSCAPE")
add_title(s9, "Infrastructure, Not Another Token Platform", width=Inches(8))

columns = ["Capability", "Simple NFT\nPlatforms", "Traditional\nERP Systems",
           "Single-Layer\nRWA Platforms", "BlockTrace"]
rows = [
    ("Hierarchical token structures", "\u2717", "\u2717", "\u2717", "\u2713"),
    ("Cryptographic provenance", "Partial", "\u2717", "Partial", "\u2713"),
    ("Enterprise ERP integration", "\u2717", "\u2713", "\u2717", "\u2713"),
    ("Recursive BOM verification", "\u2717", "\u2717", "\u2717", "\u2713"),
    ("Cross-org verifiable trust", "Partial", "\u2717", "Partial", "\u2713"),
]

table_x = Inches(0.83)
table_y = Inches(2.6)
table_w = Inches(11.67)
col_widths = [Inches(2.8)] + [Inches(2.2)] * 4
header_h = Inches(0.55)
row_h = Inches(0.45)

add_rounded_rect(s9, table_x, table_y, table_w, header_h, fill_color=ACCENT)
cx = table_x
for j, col in enumerate(columns):
    w = col_widths[j]
    add_text_box(s9, cx + Inches(0.2), table_y, w, header_h,
                 col, font_size=11, color=BLACK, bold=True,
                 alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)
    cx += w

for i, (cap, *vals) in enumerate(rows):
    ry = table_y + header_h + i * row_h
    row_fill = CARD_BG if i % 2 == 0 else DARKER_BG
    add_rounded_rect(s9, table_x, ry, table_w, row_h, fill_color=row_fill)
    cx = table_x
    for j, val in enumerate([cap] + list(vals)):
        w = col_widths[j]
        if j == 0:
            vc, fs, al = WHITE, 11, PP_ALIGN.LEFT
        elif val == "\u2713" and j == 4:
            vc, fs, al = ACCENT, 14, PP_ALIGN.CENTER
        elif val == "\u2713":
            vc, fs, al = GRAY_LIGHT, 14, PP_ALIGN.CENTER
        elif val == "Partial":
            vc, fs, al = GRAY_MED, 10, PP_ALIGN.CENTER
        else:
            vc, fs, al = GRAY_DARK, 14, PP_ALIGN.CENTER
        add_text_box(s9, cx + Inches(0.15), ry, w, row_h,
                     val, font_size=fs, color=vc,
                     bold=(val == "\u2713" and j == 4),
                     font_name="JetBrains Mono" if j > 0 else "Inter",
                     alignment=al)
        cx += w
add_slide_number(s9, 9)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10 - Business Model (Brickwork)
# ═══════════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(blank_layout)
set_slide_bg(s10)
add_section_label(s10, "BUSINESS MODEL")
add_title(s10, "Revenue Architecture", width=Inches(8))
add_description(s10,
    "Five compounding revenue streams with built-in land-and-expand mechanics. "
    "Every new asset deepens platform engagement across all streams.",
    width=Inches(7))

bm_w = Inches(3.67)
bm_h = Inches(2.7)
bm_y_top = Inches(3.0)
bm_y_bot = bm_y_top + bm_h + brick_gap

bm_top = [
    ("RECURRING", True, "*", "SaaS Subscription",
     "Tiered platform access: Starter, Growth, and Enterprise plans. "
     "Predictable ARR base that grows with seat count and org adoption. "
     "85%+ gross margins."),
    ("USAGE-BASED", False, "*", "Token Minting Fees",
     "Per-token fee for each asset or sub-token created. Revenue scales "
     "linearly with asset volume \u2014 a single enterprise can mint "
     "10K\u20131M+ tokens annually. Volume discounts drive lock-in."),
    ("USAGE-BASED", False, "*", "On-Chain Anchoring Fees",
     "Per-event fee for immutable proof anchoring. Batched for cost "
     "efficiency. High-margin revenue stream \u2014 our cost per anchor "
     "is a fraction of what customers pay."),
]
bm_bot = [
    ("HIGH-VALUE", False, "*", "Enterprise Deployments",
     "6-figure+ ACV for private cloud, on-premise, and hybrid "
     "deployments with dedicated support and SLA guarantees. Custom "
     "integrations create deep switching costs."),
    ("ADD-ON", False, "*", "Analytics & Compliance",
     "Premium add-on for asset graph analytics, automated compliance "
     "reporting, and predictive insights. Expands ACV 30\u201350% per "
     "enterprise account."),
]

for i, (badge, ba, icon, t, d) in enumerate(bm_top):
    bx = Inches(0.83) + i * (bm_w + brick_gap)
    add_card_brickwork(s10, bx, bm_y_top, bm_w, bm_h, t, d,
                       icon_text=icon, badge_text=badge,
                       badge_accent=ba)

bm_bot_start = Inches(0.83) + (3 * bm_w + 2 * brick_gap -
                                (2 * bm_w + brick_gap)) / 2
for i, (badge, ba, icon, t, d) in enumerate(bm_bot):
    bx = bm_bot_start + i * (bm_w + brick_gap)
    add_card_brickwork(s10, bx, bm_y_bot, bm_w, bm_h, t, d,
                       icon_text=icon, badge_text=badge,
                       badge_accent=ba)
add_slide_number(s10, 10)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 11 - GTM & Expansion
# ═══════════════════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(blank_layout)
set_slide_bg(s11)
add_section_label(s11, "GO-TO-MARKET")
add_title(s11, "Market Entry & Expansion Strategy", width=Inches(10))
add_description(s11,
    "Land in regulated industries where traceability is mandatory, then "
    "expand through supply chain network effects and platform stickiness.",
    width=Inches(6.5))

# Left column
left_x = Inches(0.83)
left_y = Inches(2.4)
left_w = Inches(5.8)
add_text_box(s11, left_x, left_y, Inches(3), Inches(0.25),
             "BEACHHEAD VERTICALS", font_size=9, color=ACCENT,
             bold=True, font_name="JetBrains Mono")

verticals = [
    ("Manufacturing & Industrial",
     "BOMs, quality certificates, component recall. EU Digital Product "
     "Passport mandate creates forced adoption by 2027."),
    ("Real Estate & Construction",
     "Property tokens, surveys, certificates. \u00A3300B+ UK market "
     "with fragmented, paper-heavy asset trails."),
    ("Pharma & Life Sciences",
     "Drug serialisation, cold-chain provenance, clinical trial audit "
     "trails. FDA DSCSA compliance mandatory."),
    ("Food & Agriculture",
     "Farm-to-fork traceability, batch recall, sustainability "
     "certification. EU regulation driving $2B+ in compliance spend."),
]

vert_y = left_y + Inches(0.4)
vert_h = Inches(0.8)
vert_gap = Inches(0.12)
for i, (vt, vd) in enumerate(verticals):
    vy = vert_y + i * (vert_h + vert_gap)
    add_rounded_rect(s11, left_x, vy, left_w, vert_h, fill_color=CARD_BG)
    add_text_box(s11, left_x + Inches(0.2), vy + Inches(0.1),
                 left_w - Inches(0.4), Inches(0.25),
                 vt, font_size=13, color=WHITE, bold=True)
    add_text_box(s11, left_x + Inches(0.2), vy + Inches(0.38),
                 left_w - Inches(0.4), Inches(0.4),
                 vd, font_size=10, color=GRAY_LIGHT, line_spacing=15)

# Right column
right_x = Inches(7.0)
right_w = Inches(5.8)
add_text_box(s11, right_x, left_y, Inches(3), Inches(0.25),
             "PRICING TIERS", font_size=9, color=ACCENT,
             bold=True, font_name="JetBrains Mono")

tiers = [
    ("Starter", "\u00A3500/mo",
     "Up to 1,000 tokens/mo. Single user team. Shared infrastructure. "
     "Ideal for pilots and POCs.", False),
    ("Growth", "\u00A32,500/mo",
     "Up to 25,000 tokens/mo. Multi-team access. Analytics module "
     "included. Volume anchoring discounts.", False),
    ("Enterprise", "Custom",
     "Unlimited tokens. Private infrastructure. Dedicated SLA. Custom "
     "integrations. White-glove onboarding. 6-figure ACV.", True),
]

tier_y = left_y + Inches(0.4)
tier_h = Inches(0.78)
tier_gap = Inches(0.12)
for i, (tname, tprice, tdesc, hi) in enumerate(tiers):
    ty = tier_y + i * (tier_h + tier_gap)
    shape = add_rounded_rect(s11, right_x, ty, right_w, tier_h,
                             fill_color=CARD_BG)
    if hi:
        shape.line.color.rgb = ACCENT
        shape.line.width = Pt(1)
    add_text_box(s11, right_x + Inches(0.2), ty + Inches(0.1),
                 Inches(1.3), Inches(0.25),
                 tname, font_size=13, color=WHITE, bold=True)
    add_text_box(s11, right_x + Inches(0.2), ty + Inches(0.38),
                 Inches(1.3), Inches(0.3),
                 tprice, font_size=16, color=ACCENT, bold=True)
    add_text_box(s11, right_x + Inches(1.6), ty + Inches(0.1),
                 right_w - Inches(2), tier_h - Inches(0.2),
                 tdesc, font_size=10, color=GRAY_LIGHT, line_spacing=15)

# Defensibility Moats
moat_y = tier_y + 3 * (tier_h + tier_gap) + Inches(0.15)
add_text_box(s11, right_x, moat_y, Inches(3), Inches(0.25),
             "DEFENSIBILITY MOATS", font_size=9, color=ACCENT,
             bold=True, font_name="JetBrains Mono")

moats = [
    ("*", "Data Gravity", "Cryptographic graphs are non-portable"),
    ("*", "Network Effects", "Each partner deepens the ecosystem"),
    ("*", "Regulatory Tailwinds",
     "Compliance mandates drive forced adoption"),
]

moat_card_w = Inches(1.8)
moat_card_h = Inches(1.1)
moat_card_gap = Inches(0.1)
moat_card_y = moat_y + Inches(0.35)
for i, (icon, mt, md) in enumerate(moats):
    mx = right_x + i * (moat_card_w + moat_card_gap)
    shape = add_rounded_rect(s11, mx, moat_card_y, moat_card_w,
                             moat_card_h, fill_color=DARKER_BG)
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(0.5)
    add_text_box(s11, mx + Inches(0.15), moat_card_y + Inches(0.1),
                 moat_card_w - Inches(0.3), Inches(0.25),
                 mt, font_size=10, color=WHITE, bold=True)
    add_text_box(s11, mx + Inches(0.15), moat_card_y + Inches(0.45),
                 moat_card_w - Inches(0.3), Inches(0.5),
                 md, font_size=9, color=GRAY_MED, line_spacing=14)
add_slide_number(s11, 11)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 12 - Roadmap
# ═══════════════════════════════════════════════════════════════════════
s12 = prs.slides.add_slide(blank_layout)
set_slide_bg(s12)
add_section_label(s12, "ROADMAP")
add_title(s12, "Building in Phases")

phases = [
    ("PHASE 1", True, "Core Token Engine",
     ["Hierarchical token minting", "Pack assembly & validation",
      "On-chain anchoring MVP", "REST API & basic dashboard"]),
    ("PHASE 2", True, "ERP Integrations",
     ["SAP & Oracle connectors", "Webhook event system",
      "Compliance reporting"]),
    ("PHASE 3", False, "Cross-Chain Interop",
     ["Multi-chain anchoring", "Cross-chain asset queries"]),
    ("PHASE 4", False, "AI-Driven Analytics",
     ["Predictive asset intelligence", "Automated compliance alerts"]),
]

phase_w = Inches(2.85)
phase_gap = Inches(0.2)
phase_x = Inches(0.83)
phase_y = Inches(2.4)

for i, (pname, is_active, ptitle, items) in enumerate(phases):
    px = phase_x + i * (phase_w + phase_gap)
    head_h = Inches(0.45)
    head_fill = ACCENT if i == 0 else DARKER_BG
    head_tc = BLACK if i == 0 else (ACCENT if is_active else GRAY_MED)
    add_rounded_rect(s12, px, phase_y, phase_w, head_h,
                     fill_color=head_fill)
    add_text_box(s12, px + Inches(0.2), phase_y,
                 phase_w - Inches(0.4), head_h,
                 pname, font_size=10, color=head_tc, bold=True,
                 font_name="JetBrains Mono")

    body_h = Inches(4)
    add_rounded_rect(s12, px, phase_y + head_h, phase_w, body_h,
                     fill_color=CARD_BG)
    add_text_box(s12, px + Inches(0.2), phase_y + head_h + Inches(0.15),
                 phase_w - Inches(0.4), Inches(0.35),
                 ptitle, font_size=16, color=WHITE, bold=True)

    item_color = GRAY_LIGHT if is_active else GRAY_MED
    dash_color = ACCENT if is_active else GRAY_DARK
    for j, item in enumerate(items):
        iy = phase_y + head_h + Inches(0.6) + j * Inches(0.35)
        add_text_box(s12, px + Inches(0.2), iy,
                     Inches(0.3), Inches(0.3),
                     "\u2500", font_size=10, color=dash_color,
                     font_name="JetBrains Mono")
        add_text_box(s12, px + Inches(0.45), iy,
                     phase_w - Inches(0.65), Inches(0.3),
                     item, font_size=11, color=item_color)
add_slide_number(s12, 12)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 13 - Long-Term Vision
# ═══════════════════════════════════════════════════════════════════════
s13 = prs.slides.add_slide(blank_layout)
set_slide_bg(s13)
add_section_label(s13, "LONG-TERM VISION")
add_title(s13, "Version Control for Physical Assets", width=Inches(8))
add_description(s13,
    "We\u2019re building the provenance layer for the physical world. A "
    "future where every asset has a verifiable, composable digital twin.",
    width=Inches(6.5))

vision_cards = [
    ("*", "Global Asset\nProvenance Layer",
     "A universal registry of verifiable asset histories across "
     "industries and borders."),
    ("*", "Embedded Finance\nEnablement",
     "Tokenised assets become programmable collateral for lending, "
     "insurance, and trade finance."),
    ("*", "Automated\nCompliance",
     "Regulatory reporting generated automatically from verifiable "
     "asset data."),
    ("*", "Cross-Industry\nComposability",
     "Token packs from one industry interoperate with token packs "
     "from another."),
]

vis_w = Inches(2.85)
vis_h = Inches(3.0)
vis_gap = Inches(0.2)
vis_y = Inches(3.6)

for i, (icon, vt, vd) in enumerate(vision_cards):
    vx = Inches(0.83) + i * (vis_w + vis_gap)
    shape = add_rounded_rect(s13, vx, vis_y, vis_w, vis_h,
                             fill_color=CARD_BG)
    if i == 3:
        shape.line.color.rgb = ACCENT
        shape.line.width = Pt(1)
    add_text_box(s13, vx, vis_y + Inches(0.35), vis_w, Inches(0.5),
                 icon, font_size=30, color=ACCENT,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(s13, vx + Inches(0.2), vis_y + Inches(0.9),
                 vis_w - Inches(0.4), Inches(0.6),
                 vt, font_size=14, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER, line_spacing=20)
    add_text_box(s13, vx + Inches(0.2), vis_y + Inches(1.6),
                 vis_w - Inches(0.4), Inches(1.0),
                 vd, font_size=11, color=GRAY_LIGHT,
                 alignment=PP_ALIGN.CENTER, line_spacing=18)
add_slide_number(s13, 13)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 14 - Closing
# ═══════════════════════════════════════════════════════════════════════
s14 = prs.slides.add_slide(blank_layout)
set_slide_bg(s14)
add_accent_line(s14, Inches(0), Inches(0), SLIDE_W, Pt(3))
add_text_box(s14, Inches(2.5), Inches(2.5), Inches(8.333), Inches(1),
             "\u201CBuild Trust Into Your Assets.\u201D",
             font_size=40, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)
add_accent_line(s14, Inches(6), Inches(3.7), Inches(1.333), Pt(2))
add_text_box(s14, Inches(2.5), Inches(4.0), Inches(8.333), Inches(0.6),
             "BlockTrace", font_size=20, color=ACCENT, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(s14, Inches(2.5), Inches(4.7), Inches(8.333), Inches(0.35),
             "contact@blocktrace.io", font_size=12, color=GRAY_MED,
             font_name="JetBrains Mono", alignment=PP_ALIGN.CENTER)
add_text_box(s14, Inches(2.5), Inches(5.05), Inches(8.333), Inches(0.35),
             "www.blocktrace.io", font_size=12, color=GRAY_MED,
             font_name="JetBrains Mono", alignment=PP_ALIGN.CENTER)
add_slide_number(s14, 14)

# Save
output_path = "BlockTrace_Pitch_Deck.pptx"
prs.save(output_path)
print(f"\u2705 Saved {output_path} \u2014 {len(prs.slides)} slides")
